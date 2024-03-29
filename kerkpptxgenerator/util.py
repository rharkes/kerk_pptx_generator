"""
Tools to make pptx
For testing with IPython:
%load_ext autoreload
%autoreload 2
"""
from pathlib import Path
import re
from typing import List

from pptx import Presentation
from pptx.util import Cm as Centimeter
from PIL.PngImagePlugin import PngImageFile


class SlideProperties:
    def __init__(
        self,
        topmargin: float = 0.5,
        leftmargin: float = 0.5,
        width: float = 16,
        height: float = 9,
    ):
        self.topmargin = Centimeter(topmargin)
        self.leftmargin = Centimeter(leftmargin)
        self.width = Centimeter(width)
        self.height = Centimeter(height)
        self.availablewidth = 0.0
        self.availableheight = 0.0
        self.ratio = 0.0
        self.recalculate()

    def recalculate(self) -> None:
        self.availablewidth = self.width - 2 * self.leftmargin
        self.availableheight = self.height - 2 * self.topmargin
        self.ratio = self.availablewidth / self.availableheight

    def setratio(self, newratio: float) -> None:
        if newratio < self.ratio:  # width must decrease by increasing leftmargin
            requiredwidth = (newratio / self.ratio) * self.availablewidth
            self.leftmargin += (self.availablewidth - requiredwidth) / 2
            self.recalculate()
        else:  # height must decrease directly, topmargin stays the same
            self.availableheight *= self.ratio / newratio
            self.ratio = self.availablewidth / self.availableheight


class SongList:
    """
    Accepts the image path and list path and generates paths to images.
    """

    def __init__(self, img_pth: Path, list_pth: Path):
        self.img_pth = Path(img_pth)
        self.list_pth = Path(list_pth)
        self.idx = -1  # type: int
        self.paths = self.getpaths()

    def getpaths(self) -> List[Path]:
        with open(self.list_pth) as f:
            lines = f.readlines()
        allpaths = []
        for line in lines:
            values = line.split(" ", 1)  # splits het lied af van de coupletten
            if len(values) == 1:  # alle coupletten
                iml = [
                    x
                    for x in self.img_pth.glob(f"projectie-{values[0].strip()}-muziek*")
                ]
                if len(iml) == 0:
                    print(f"WAARSCHUWING: Lied {values[0].strip()} niet gevonden.")
                else:
                    # print(f"Lied{values[0]} : {iml}")
                    pass
            else:  # coupletten gespecificeerd
                iml = []
                coupletten = values[1].split(",")
                if (
                    values[0] == "1004"
                ):  # Lied 1004 is anders, heeft een opening en sluiting, coupletten staan ertussen
                    iml += [x for x in self.img_pth.glob("projectie-1004-muziek-1.png")]
                    for c in coupletten:
                        csi = int(c.strip())
                        iml += [
                            x
                            for x in self.img_pth.glob(
                                f"projectie-1004-muziek-{csi+1}.png"
                            )
                        ]
                    iml += [x for x in self.img_pth.glob("projectie-1004-muziek-9.png")]
                    coupletten = []
                for c in coupletten:
                    cs = c.strip()
                    new_iml = [
                        x
                        for x in self.img_pth.glob(
                            f"projectie-{values[0]}-muziek-couplet-{cs}*"
                        )
                    ]
                    if len(new_iml) == 0:
                        print(
                            f"WAARSCHUWING: Lied {values[0]} couplet {cs} niet gevonden."
                        )
                    else:
                        # print(f"Lied{values[0]} couplet{cs} : {new_iml}")
                        pass
                    iml += new_iml
            allpaths += iml
        return allpaths


def make_presentation(slidecfg: dict[str, float]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Centimeter(slidecfg["width"])
    prs.slide_height = Centimeter(slidecfg["height"])
    return prs


def add_pictureslide(
    prs: Presentation, img_path: Path, cfg: dict[str, float]
) -> Presentation:
    """
    Add a slide with a picture from img_path with specific margin
    :param prs:
    :param img_path:
    :param cfg:
    :return:
    """
    sp = SlideProperties(
        cfg["topmargin"], cfg["leftmargin"], cfg["width"], cfg["height"]
    )
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    img_path, imrat = crop_picture(img_path)
    sp.setratio(imrat)
    slide.shapes.add_picture(
        str(img_path),
        left=sp.leftmargin,
        top=sp.topmargin,
        width=sp.availablewidth,
        height=sp.availableheight,
    )
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if cfg["include_notes"]:
        filename = img_path.name
        nrs = re.findall("[0-9]+", filename)
        if str.__contains__(filename, "couplet"):
            text_frame.text = f"Lied {nrs[0]} couplet {nrs[1]}"
        else:
            text_frame.text = f"Lied {nrs[0]}"
    return prs


def crop_picture(img_path_in: Path) -> tuple[Path, float]:
    """
    Crop picture to remove all white space around it
    Save the red channel as grayscale
    :param img_path_in:
    :return:
    """
    img_path_out = Path(img_path_in.parent, "crops", img_path_in.stem + "_crp.png")
    if img_path_out.exists():
        im = PngImageFile(img_path_out)
        return img_path_out, im.width / im.height
    """
    Python only evaluates the portion of a logical expression as is necessary to determine the outcome, 
    and returns the last value examined as the result of the expression. 
    So if the expression above is false (0), Python does not look at the second operand, 
    and thus returns 0. Otherwise, it returns 255.
    """
    im = PngImageFile(img_path_in)
    source = im.split()
    mask = source[0].point(lambda i: i < 255 and 255)
    bbox = mask.getbbox()
    img_crp = source[0].crop(bbox)
    img_path_out.parent.mkdir(exist_ok=True)
    img_crp.save(img_path_out)
    return img_path_out, img_crp.width / img_crp.height
